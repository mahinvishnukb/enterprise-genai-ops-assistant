"""
Document loaders: turn an uploaded file's bytes into plain text, regardless
of source format. Each parser is imported lazily inside its function so the
module still loads (and the other formats still work) even if one optional
dependency isn't installed.
"""
from pathlib import Path


def extract_text(file_path: str) -> str:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(file_path)
    if suffix == ".docx":
        return _extract_docx(file_path)
    if suffix == ".pptx":
        return _extract_pptx(file_path)
    if suffix == ".csv":
        return _extract_csv(file_path)
    # .txt, .md, and unknown types: read as plain text.
    return Path(file_path).read_text(errors="ignore")


def _extract_pdf(file_path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(file_path: str) -> str:
    import docx

    document = docx.Document(file_path)
    return "\n".join(p.text for p in document.paragraphs)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _extract_csv(file_path: str) -> str:
    import csv

    rows = []
    with open(file_path, newline="", errors="ignore") as f:
        for row in csv.reader(f):
            rows.append(", ".join(row))
    return "\n".join(rows)
